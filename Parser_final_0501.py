import re
import os
import shutil
import tempfile
import extract_msg
import logging
import datetime
import json
from PyPDF2 import PdfReader  # pip install PyPDF2
from docx import Document      # pip install python-docx
from pptx import Presentation  # pip install python-pptx
from openpyxl import load_workbook # pip install openpyxl

# Module-level logger
logger = logging.getLogger(__name__)


def format_email_body(body: str) -> str:
    """
    Normalize and clean the email body:
      - Convert CRLF/CR to LF
      - Collapse multiple spaces/tabs
      - Collapse multiple blank lines
      - Trim leading/trailing whitespace
    """
    logger.debug("Formatting email body (raw length=%d)", len(body))
    body = body.replace('\r\n', '\n').replace('\r', '\n')
    body = re.sub(r'[ \t]+', ' ', body)
    body = re.sub(r'\n\s*\n', '\n\n', body)
    formatted = body.strip()
    logger.debug("Formatted email body (length=%d)", len(formatted))
    return formatted


def safe_json_load(raw):
    """Try loading a string into JSON; fallback to original string"""
    try:
        return json.loads(raw)
    except (json.JSONDecodeError, TypeError) as e:
        logger.debug("safe_json_load failed: %s", e)
        return raw

def parse_analysis_field(data: dict) -> dict:
    """
    Extracts structured fields from LLM response in either of these formats:
    
    - Inline-numbered:
      '1. Result: ... 2. Category: ... 3. Explanation: ... 4. Citation: ...'
    
    - Newline-separated:
      'Result: ...\nCategory: ...\nExplanation: ...\nCitation: ...'
    """
    answer = data.get('answer')
    if not isinstance(answer, str):
        return {}

    parsed = {}

    # Try numbered format first: 1. Result: ... 2. Category: ...
    numbered_matches = re.findall(r"\d+\.\s*([\w\s]+):\s*(.*?)(?=\d+\.|$)", answer, re.DOTALL)

    if numbered_matches:
        for key, value in numbered_matches:
            key_clean = key.strip().lower().replace(' ', '_')
            parsed[key_clean] = value.strip()
            logger.debug("Parsed field '%s': '%s'", key_clean, parsed[key_clean])
        return parsed

    # Fallback: plain "Key: value" format with newlines
    line_matches = re.findall(r"([\w\s]+):\s*(.*?)(?=\n[\w\s]+:|$)", answer, re.DOTALL)
    for key, value in line_matches:
        key_clean = key.strip().lower().replace(' ', '_')
        parsed[key_clean] = value.strip()
        logger.debug("Parsed field '%s': '%s'", key_clean, parsed[key_clean])

    return parsed


def extract_pdf_text(pdf_path: str) -> str:
    """Extract text from a PDF file"""
    logger.debug("Extracting text from PDF: %s", pdf_path)
    reader = PdfReader(pdf_path)
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text)
    combined = "\n".join(texts).strip()
    logger.debug("Extracted PDF text length=%d", len(combined))
    return combined


def extract_docx_text(docx_path: str) -> str:
    """Extract text from a .docx Word document"""
    logger.debug("Extracting text from DOCX: %s", docx_path)
    doc = Document(docx_path)
    paragraphs = [para.text for para in doc.paragraphs if para.text]
    combined = "\n".join(paragraphs).strip()
    logger.debug("Extracted DOCX text length=%d", len(combined))
    return combined


def extract_pptx_text(pptx_path: str) -> str:
    """Extract text from a .pptx PowerPoint presentation"""
    logger.debug("Extracting text from PPTX: %s", pptx_path)
    prs = Presentation(pptx_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text)
    combined = "\n".join(texts).strip()
    logger.debug("Extracted PPTX text length=%d", len(combined))
    return combined


def extract_excel_text(xlsx_path: str) -> str:
    """Extract text from an Excel workbook (.xlsx) by reading all cells"""
    logger.debug("Extracting text from Excel: %s", xlsx_path)
    wb = load_workbook(xlsx_path, read_only=True)
    texts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell) for cell in row if cell is not None]
            if row_text:
                texts.append("\t".join(row_text))
    combined = "\n".join(texts).strip()
    logger.debug("Extracted Excel text length=%d", len(combined))
    return combined


def parse_msg(file_path: str) -> dict:
    """
    Parse a .msg file:
      1) Extract metadata and format the main body.
      2) Detect embedded attachments (MSG, PDF, DOCX, PPTX, XLSX), extract text.
      3) Return dict with metadata, body, pdf_attachments, other attachments, nested_emails.
    """
    logger.debug("Starting parse_msg for file: %s", file_path)
    msg = extract_msg.Message(file_path)
    try:
        # Metadata
        metadata = {
            "From": msg.sender,
            "To": msg.to,
            "Cc": msg.cc,
            "Bcc": "",
            "Date": msg.date
        }
        # Date normalization
        date_val = metadata["Date"]
        if isinstance(date_val, datetime.datetime):
            metadata["Date"] = date_val.isoformat()
        else:
            metadata["Date"] = str(date_val)

        # Body
        raw_body = msg.body or getattr(msg, "htmlBody", "") or ""
        body = format_email_body(raw_body) if raw_body else ""

        result = {"metadata": metadata, "body": body}
        attachments = msg.attachments or []
        logger.debug("Total attachments found: %d", len(attachments))

        # Lists to collect extracted contents
        pdf_list = []
        docx_list = []
        pptx_list = []
        xlsx_list = []
        nested_list = []

        for att in attachments:
            filename = att.getFilename() or ""
            logger.debug("Inspecting attachment: %s", filename)

            # PDF
            if filename.lower().endswith('.pdf'):
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmpf:
                    tmpf.write(att.data)
                    tmp_pdf = tmpf.name
                try:
                    text = extract_pdf_text(tmp_pdf)
                    pdf_list.append({"filename": filename, "content": text})
                except Exception as e:
                    logger.error("Error extracting PDF %s: %s", filename, e)
                finally:
                    os.remove(tmp_pdf)

            # DOCX
            elif filename.lower().endswith('.docx'):
                with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmpf:
                    tmpf.write(att.data)
                    tmp_docx = tmpf.name
                try:
                    text = extract_docx_text(tmp_docx)
                    docx_list.append({"filename": filename, "content": text})
                except Exception as e:
                    logger.error("Error extracting DOCX %s: %s", filename, e)
                finally:
                    os.remove(tmp_docx)

            # PPTX
            elif filename.lower().endswith('.pptx'):
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmpf:
                    tmpf.write(att.data)
                    tmp_pptx = tmpf.name
                try:
                    text = extract_pptx_text(tmp_pptx)
                    pptx_list.append({"filename": filename, "content": text})
                except Exception as e:
                    logger.error("Error extracting PPTX %s: %s", filename, e)
                finally:
                    os.remove(tmp_pptx)

            # XLSX
            elif filename.lower().endswith(('.xlsx', '.xls')):
                with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as tmpf:
                    tmpf.write(att.data)
                    tmp_xlsx = tmpf.name
                try:
                    text = extract_excel_text(tmp_xlsx)
                    xlsx_list.append({"filename": filename, "content": text})
                except Exception as e:
                    logger.error("Error extracting Excel %s: %s", filename, e)
                finally:
                    os.remove(tmp_xlsx)

            # Nested MSG
            elif filename.lower().endswith('.msg'):
                tmpdir = tempfile.mkdtemp()
                try:
                    save_type, paths = att.save(customPath=tmpdir, extractEmbedded=True)
                    if isinstance(paths, str):
                        paths = [paths]
                    for p in paths:
                        nested_data = parse_msg(p)
                        nested_list.append(nested_data)
                except Exception as e:
                    logger.error("Error handling nested .msg %s: %s", filename, e)
                finally:
                    shutil.rmtree(tmpdir, ignore_errors=True)

        if pdf_list:
            result['pdf_attachments'] = pdf_list
        if docx_list:
            result['docx_attachments'] = docx_list
        if pptx_list:
            result['pptx_attachments'] = pptx_list
        if xlsx_list:
            result['excel_attachments'] = xlsx_list
        if nested_list:
            result['nested_emails'] = nested_list

        return result

    finally:
        try:
            msg.close()
        except:
            pass

def parse_email(file_path: str) -> dict:
    logger.debug("In parse_email with file: %s", file_path)
    if file_path.lower().endswith('.msg'):
        return parse_msg(file_path)
    error_msg = "Unsupported file format. Only .msg files are supported."
    logger.error(error_msg)
    raise ValueError(error_msg)
