import re
import os
import shutil
import tempfile
import extract_msg
import logging
import datetime
import json
from PyPDF2 import PdfReader  # pip install PyPDF2
from docx import Document      # pip install python-docx
from pptx import Presentation  # pip install python-pptx
from openpyxl import load_workbook # pip install openpyxl

# Module-level logger
logger = logging.getLogger(__name__)

def format_email_body(body: str) -> str:
    logger.debug("Formatting email body (raw length=%d)", len(body))
    body = body.replace('\r\n', '\n').replace('\r', '\n')
    body = re.sub(r'[ \t]+', ' ', body)
    body = re.sub(r'\n\s*\n', '\n\n', body)
    formatted = body.strip()
    logger.debug("Formatted email body (length=%d)", len(formatted))
    return formatted

def safe_json_load(raw):
    try:
        return json.loads(raw)
    except (json.JSONDecodeError, TypeError) as e:
        logger.debug("safe_json_load failed: %s", e)
        return raw

def parse_analysis_field(data: dict) -> dict:
    """
    Assumes the LLM response is already in structured JSON format:
    {
        "classification": "...",
        "category": "...",
        "explanation": "...",
        "citation": "..."
    }
    """
    expected_fields = ["classification", "category", "explanation", "citation"]
    parsed = {}
    for field in expected_fields:
        value = data.get(field)
        if value:
            parsed[field] = value.strip() if isinstance(value, str) else value
            logger.debug("Parsed JSON field '%s': '%s'", field, parsed[field])
    return parsed

def extract_pdf_text(pdf_path: str) -> str:
    logger.debug("Extracting text from PDF: %s", pdf_path)
    reader = PdfReader(pdf_path)
    texts = [page.extract_text() for page in reader.pages if page.extract_text()]
    combined = "\n".join(texts).strip()
    logger.debug("Extracted PDF text length=%d", len(combined))
    return combined

def extract_docx_text(docx_path: str) -> str:
    logger.debug("Extracting text from DOCX: %s", docx_path)
    doc = Document(docx_path)
    paragraphs = [para.text for para in doc.paragraphs if para.text]
    combined = "\n".join(paragraphs).strip()
    logger.debug("Extracted DOCX text length=%d", len(combined))
    return combined

def extract_pptx_text(pptx_path: str) -> str:
    logger.debug("Extracting text from PPTX: %s", pptx_path)
    prs = Presentation(pptx_path)
    texts = [shape.text.strip() for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
    combined = "\n".join(texts).strip()
    logger.debug("Extracted PPTX text length=%d", len(combined))
    return combined

def extract_excel_text(xlsx_path: str) -> str:
    logger.debug("Extracting text from Excel: %s", xlsx_path)
    wb = load_workbook(xlsx_path, read_only=True)
    texts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell) for cell in row if cell is not None]
            if row_text:
                texts.append("\t".join(row_text))
    combined = "\n".join(texts).strip()
    logger.debug("Extracted Excel text length=%d", len(combined))
    return combined

def parse_msg(file_path: str) -> dict:
    logger.debug("Starting parse_msg for file: %s", file_path)
    msg = extract_msg.Message(file_path)
    try:
        metadata = {
            "From": msg.sender,
            "To": msg.to,
            "Cc": msg.cc,
            "Bcc": "",
            "Date": msg.date
        }
        date_val = metadata["Date"]
        metadata["Date"] = date_val.isoformat() if isinstance(date_val, datetime.datetime) else str(date_val)

        raw_body = msg.body or getattr(msg, "htmlBody", "") or ""
        body = format_email_body(raw_body) if raw_body else ""
        result = {"metadata": metadata, "body": body}

        attachments = msg.attachments or []
        logger.debug("Total attachments found: %d", len(attachments))

        pdf_list, docx_list, pptx_list, xlsx_list, nested_list = [], [], [], [], []

        for att in attachments:
            filename = att.getFilename() or ""
            logger.debug("Inspecting attachment: %s", filename)

            suffix = os.path.splitext(filename)[1].lower()
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmpf:
                tmpf.write(att.data)
                tmp_path = tmpf.name

            try:
                if suffix == '.pdf':
                    text = extract_pdf_text(tmp_path)
                    pdf_list.append({"filename": filename, "content": text})
                elif suffix == '.docx':
                    text = extract_docx_text(tmp_path)
                    docx_list.append({"filename": filename, "content": text})
                elif suffix == '.pptx':
                    text = extract_pptx_text(tmp_path)
                    pptx_list.append({"filename": filename, "content": text})
                elif suffix in ['.xlsx', '.xls']:
                    text = extract_excel_text(tmp_path)
                    xlsx_list.append({"filename": filename, "content": text})
                elif suffix == '.msg':
                    tmpdir = tempfile.mkdtemp()
                    try:
                        save_type, paths = att.save(customPath=tmpdir, extractEmbedded=True)
                        for p in ([paths] if isinstance(paths, str) else paths):
                            nested_data = parse_msg(p)
                            nested_list.append(nested_data)
                    finally:
                        shutil.rmtree(tmpdir, ignore_errors=True)
            except Exception as e:
                logger.error("Error extracting %s: %s", filename, e)
            finally:
                os.remove(tmp_path)

        if pdf_list: result['pdf_attachments'] = pdf_list
        if docx_list: result['docx_attachments'] = docx_list
        if pptx_list: result['pptx_attachments'] = pptx_list
        if xlsx_list: result['excel_attachments'] = xlsx_list
        if nested_list: result['nested_emails'] = nested_list

        return result
    finally:
        try:
            msg.close()
        except:
            pass

def parse_email(file_path: str) -> dict:
    logger.debug("In parse_email with file: %s", file_path)
    if file_path.lower().endswith('.msg'):
        return parse_msg(file_path)
    error_msg = "Unsupported file format. Only .msg files are supported."
    logger.error(error_msg)
    raise ValueError(error_msg)
